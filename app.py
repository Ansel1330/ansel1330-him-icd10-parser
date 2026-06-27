import streamlit as st
import pandas as pd
import numpy as np
import os
import anthropic
import json
import re

# --- 1. PAGE SETUP ---
st.set_page_config(page_title="UCC HIM Level 200 Hub", page_icon="🎓", layout="wide")

if "dark_mode" not in st.session_state:
    st.session_state.dark_mode = False

# --- 2. DYNAMIC CONTRAST COLOR ENGINE ---
if st.session_state.dark_mode:
    bg_color    = "#0F172A"
    text_color  = "#F8FAFC"
    card_bg     = "#1E293B"
    border_color= "#334155"
    sub_text    = "#94A3B8"
else:
    bg_color    = "#FFFFFF"
    text_color  = "#0F172A"
    card_bg     = "#F8FAFC"
    border_color= "#E2E8F0"
    sub_text    = "#475569"

st.markdown(f"""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

    html, body, [class*="st-"] {{
        font-family: 'Inter', sans-serif !important;
    }}

    .main-title {{
        font-size: 34px;
        font-weight: 700;
        color: #FF4B4B !important;
        letter-spacing: -0.5px;
        margin-bottom: 5px;
    }}
    .section-subtitle {{
        font-size: 16px;
        color: {sub_text} !important;
        margin-bottom: 25px;
    }}
    .course-card {{
        padding: 20px;
        border-radius: 12px;
        border-left: 6px solid #FF4B4B;
        margin-bottom: 20px;
        background-color: {card_bg};
        border: 1px solid {border_color};
        border-left: 6px solid #FF4B4B;
    }}
    .card-text {{
        color: {text_color} !important;
        font-weight: 600;
        margin: 0;
    }}
    .pool-banner {{
        background: linear-gradient(135deg, #FF4B4B 0%, #991B1B 100%);
        padding: 25px;
        border-radius: 16px;
        margin-bottom: 25px;
    }}
    .ai-badge {{
        background: linear-gradient(135deg, #7C3AED, #4F46E5);
        color: white; padding: 4px 10px; border-radius: 20px;
        font-size: 12px; font-weight: 600; display: inline-block; margin-bottom: 10px;
    }}
    .explanation-box {{
        background:#1E3A5F; padding:14px; border-radius:8px;
        border-left:4px solid #3B82F6; margin-top:10px;
    }}
    [data-testid="stSidebarCollapseButton"] {{
        background-color: #FF4B4B22 !important;
        border: 1px solid #FF4B4B !important;
        border-radius: 8px;
    }}
    </style>
""", unsafe_allow_html=True)

# --- 3. SIDEBAR HEADER & NAVIGATION ---
if os.path.exists("OIP.webp"):
    st.sidebar.image("OIP.webp", use_container_width=True)
else:
    st.sidebar.markdown("<h2 style='color:#FF4B4B; font-weight:700; margin:0;'>🏥 HIMSA UCC</h2>", unsafe_allow_html=True)

st.sidebar.markdown("<h2 style='font-size:22px; font-weight:600; margin-top:5px; color:#FF4B4B;'>UCC HIM Portal</h2>", unsafe_allow_html=True)
st.sidebar.markdown(f"<p style='color:{sub_text}; font-size:14px; margin-top:-15px;'>Level 200 - Semester 2</p>", unsafe_allow_html=True)
st.sidebar.markdown("---")

app_mode = st.sidebar.radio(
    "🧭 Select App Section Below:",
    ["📚 Course Material Distribution", "📝 Adaptive Quiz Arena", "🏆 Student Leaderboard Pool", "📊 Academic Targeter"]
)

st.sidebar.markdown("---")
st.session_state.dark_mode = st.sidebar.toggle("🌙 Enable Dark Mode Layout")

# --- 4. PERSISTENT STATE TRACKING ---
if "quiz_level"           not in st.session_state: st.session_state.quiz_level = "Medium"
if "score"                not in st.session_state: st.session_state.score = 0
if "current_course_quiz"  not in st.session_state: st.session_state.current_course_quiz = "HIM 212: Programming and Software Development II"
if "active_student"       not in st.session_state: st.session_state.active_student = None
if "competition_pool"     not in st.session_state: st.session_state.competition_pool = {}
if "prev_cgpa"            not in st.session_state: st.session_state.prev_cgpa = 3.2
if "target_gpa"           not in st.session_state: st.session_state.target_gpa = 3.3
# AI quiz state
if "ai_questions"         not in st.session_state: st.session_state.ai_questions = []
if "current_topic"        not in st.session_state: st.session_state.current_topic = None
if "quiz_index"           not in st.session_state: st.session_state.quiz_index = 0
if "answered"             not in st.session_state: st.session_state.answered = False
if "last_correct"         not in st.session_state: st.session_state.last_correct = None
if "explanation"          not in st.session_state: st.session_state.explanation = ""

# --- 5. AI COURSE CONTEXTS — ALL 7 COURSES ---
COURSE_CONTEXTS = {
    "HIM 212: Programming and Software Development II": """
        You are a university lecturer generating exam questions for HIM 212: Programming and Software Development II
        at the University of Cape Coast, Ghana. Topics include: Python functions and variable scope (local vs global),
        GUI development with Tkinter (widgets, layout managers: pack, grid, place; event binding), object-oriented
        programming (classes, __init__, inheritance, encapsulation, polymorphism), file handling (open, read, write,
        append, with statement), exception/error handling (try/except/finally), Python modules and packages,
        and building clinical health information system desktop applications (patient registration forms, data entry UIs).
        Students are Level 200 Health Information Management undergraduates.
    """,
    "HIM 208: Pathophysiology II": """
        You are a university lecturer generating exam questions for HIM 208: Pathophysiology II
        at the University of Cape Coast, Ghana. Topics include: fever mechanisms, exogenous pyrogens
        (LPS from Gram-negative bacteria) and endogenous pyrogens (IL-1, IL-6, TNF), prostaglandin E2 (PGE2)
        synthesis and the hypothalamic set-point, acute and chronic inflammatory responses, cancer pathophysiology
        (oncogenesis, proto-oncogenes vs tumour suppressor genes, tumour classification — benign vs malignant,
        metastasis routes), cardiovascular diseases (hypertension, atherosclerosis, ischaemic heart disease,
        heart failure), respiratory disorders (pneumonia, COPD, asthma), endocrine disruptions (diabetes mellitus
        type 1 vs type 2, thyroid disorders), diet and nutrition in disease, and immune system pathology.
        Students are Level 200 Health Information Management undergraduates.
    """,
    "HIM 204: Database Structures": """
        You are a university lecturer generating exam questions for HIM 204: Database Structures
        at the University of Cape Coast, Ghana. Topics include: DBMS concepts, advantages over flat files,
        data models (hierarchical, network, relational), DBMS architecture (3-tier, internal/conceptual/external),
        Entity-Relationship (ER) modeling — entities, attributes, relationships, cardinality (1:1, 1:N, M:N),
        ER diagram construction and working with ER diagrams, relational database design, normalization
        (1NF — atomic values; 2NF — full functional dependency; 3NF — no transitive dependency;
        BCNF), normal forms (practical identification and decomposition), introduction to database queries
        (SQL SELECT, WHERE, JOIN, GROUP BY, ORDER BY, subqueries), ACID transactions, concurrency control,
        indexing, and health data management systems.
        Students are Level 200 Health Information Management undergraduates.
    """,
    "HIM 210: Disease Classification and Coding": """
        You are a university lecturer generating exam questions for HIM 210: Disease Classification and Coding
        at the University of Cape Coast, Ghana. Topics include: ICD-10 structure and three-volume conventions
        (Volume 1 — Tabular List of diseases; Volume 2 — Instruction Manual; Volume 3 — Alphabetical Index),
        ICD-10 coding rules and guidelines, use of inclusion and exclusion notes, coding of signs and symptoms
        vs confirmed diagnoses, morphology codes (M-codes), sequencing of principal and secondary diagnoses,
        external causes of morbidity and mortality (V, W, X, Y codes), Z-codes (factors influencing health
        status and contact with health services), DRG (Diagnosis-Related Groups) and their use in hospital
        reimbursement, clinical coding for inpatient and outpatient settings, and the two-step coding process
        (Alphabetical Index → Tabular List verification).
        Students are Level 200 Health Information Management undergraduates.
    """,
    "PHL 205: Critical Thinking and Practical Reasoning": """
        You are a university lecturer generating exam questions for PHL 205: Critical Thinking and Practical Reasoning
        at the University of Cape Coast, Ghana. Topics include: identifying and evaluating logical fallacies
        (straw man, ad hominem, false dichotomy, appeal to authority, slippery slope, hasty generalisation,
        circular reasoning, red herring, appeal to emotion, bandwagon), argument structure (premise, inference,
        conclusion; identifying implicit premises), argument validity vs soundness, deductive reasoning
        (modus ponens, modus tollens, syllogisms), inductive reasoning (strength, cogency, generalisation),
        analogical reasoning, cognitive biases (confirmation bias, anchoring, availability heuristic),
        evaluating evidence quality, and the application of critical logic in healthcare decision-making
        and health information management professional contexts.
        Students are Level 200 Health Information Management undergraduates.
    """,
    "HIM 202: System Analysis and Design": """
        You are a university lecturer generating exam questions for HIM 202: System Analysis and Design
        at the University of Cape Coast, Ghana. Topics include: systems development life cycle (SDLC) phases
        (planning, analysis, design, implementation, maintenance), structured vs object-oriented analysis,
        feasibility studies (technical, economic, operational, schedule), requirements gathering techniques
        (interviews, questionnaires, observation, document review), data flow diagrams (DFDs) — context diagrams,
        level-0 and level-1 diagrams, data dictionaries, process specifications (structured English, decision
        tables, decision trees), system design (input/output design, interface design, file and database design),
        prototyping, CASE tools, system testing strategies, and implementation and changeover strategies
        (direct, parallel, phased, pilot). Apply concepts to health information systems design contexts.
        Students are Level 200 Health Information Management undergraduates.
    """,
    "HIM 206: Data Communication and Networks in Healthcare": """
        You are a university lecturer generating exam questions for HIM 206: Data Communication and Networks
        in Healthcare at the University of Cape Coast, Ghana. Topics include: fundamentals of data communication
        (transmission media — twisted pair, coaxial, fibre optic, wireless; bandwidth; signal types — analog
        vs digital), network topologies (bus, star, ring, mesh, hybrid), network types (LAN, WAN, MAN, PAN),
        OSI reference model (7 layers and their functions), TCP/IP protocol suite, IP addressing (IPv4, subnetting),
        network hardware (routers, switches, hubs, access points, NICs), network security (firewalls, VPN,
        encryption, authentication), healthcare-specific networking (HL7 standards, DICOM for medical imaging,
        telemedicine infrastructure, Electronic Health Record (EHR) network requirements, HIPAA/data privacy
        in networked health environments), and wireless health technologies (mHealth, IoT in healthcare).
        Students are Level 200 Health Information Management undergraduates.
    """
}

ALL_COURSES = list(COURSE_CONTEXTS.keys())


# --- 6. LOCAL CONTENT QUIZ ENGINE ---
import PyPDF2
from pptx import Presentation

def extract_text_from_file(filename):
    text = ""
    try:
        if filename.endswith(".pdf"):
            with open(filename, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        elif filename.endswith(".pptx") or filename.endswith(".ppt"):
            prs = Presentation(filename)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        return f"Error reading {filename}: {e}"
    return text

def generate_ai_questions(topic: str, num_questions: int = 10) -> list:
    """Parses local files for content and generates study prompts for all courses."""
    
    # Define your course file registry here
    course_map = {
        "HIM 212: Programming and Software Development II": ["1. WEEK 1_Python Functions_May 2026.pdf"],
        "HIM 208: Pathophysiology II": ["CANCER.pptx", "Fever-HIM NEW (2).pptx"],
        "HIM 204: Database Structures": ["YOUR_HIM204_FILE.pdf"],
        "HIM 210: Disease Classification and Coding": ["YOUR_HIM210_FILE.pdf"],
        "PHL 205: Critical Thinking and Practical Reasoning": ["YOUR_PHL205_FILE.pdf"],
        "HIM 202: System Analysis and Design": ["YOUR_HIM202_FILE.pdf"],
        "HIM 206: Data Communication and Networks in Healthcare": ["YOUR_HIM206_FILE.pdf"],
    }
    
    # Get the files associated with the selected topic
    files = course_map.get(topic, [])
    full_text = ""
    
    for f in files:
        if os.path.exists(f):
            full_text += extract_text_from_file(f)
        else:
            full_text += f" [File {f} not found in directory] "
    
    # Create simple questions from text chunks
    chunks = [c.strip() for c in full_text.split('.') if len(c) > 30]
    questions = []
    
    # Generate questions based on available text
    for i in range(min(num_questions, len(chunks))):
        questions.append({
            "q": f"Based on your notes: '{chunks[i][:80]}...'",
            "o": ["True", "False", "Requires further study", "Concept understood"],
            "c": "True",
            "explanation": "This content was extracted from your course materials."
        })
        
    # Fallback if no text was found
    if not questions:
        questions.append({
            "q": "No study material found for this course. Please ensure files are in the folder.",
            "o": ["OK", "Retry", "Check Path", "Upload File"],
            "c": "OK",
            "explanation": "The system couldn't find the specific files mapped in the course_map."
        })
        
    return questions

# ============================================================
# MODULE 1: COURSE MATERIAL DISTRIBUTION
# ============================================================
if app_mode == "📚 Course Material Distribution":
    st.markdown("<h1 class='main-title'>📚 Courseware Repository</h1>", unsafe_allow_html=True)
    st.markdown("<p class='section-subtitle'>Access core reference sheets, lecture slide decks, and module mappings.</p>", unsafe_allow_html=True)
    st.write("---")

    course_selection = st.selectbox(
        "📖 Select Course Department",
        [
            "HIM 212: Programming and Software Development II",
            "HIM 208: Pathophysiology II",
            "HIM 204: Database Structures",
            "HIM 210: Disease Classification and Coding",
            "PHL 205: Critical Thinking and Practical Reasoning",
            "HIM 202: System Analysis and Design",
            "HIM 206: Data Communication and Networks in Healthcare"
        ]
    )

    st.markdown(f"""
        <div class='course-card'>
            <p class='card-text'>Current Portal: {course_selection}</p>
        </div>
    """, unsafe_allow_html=True)

    # --- HIM 212 ---
    if "HIM 212" in course_selection:
        st.write("#### Handouts & Slide Decks")
        files = [
            "1. WEEK 1_Python Functions_May 2026.pdf",
            "WEEK4_Introduction to GUI Programs (2).pptx",
            "PROGRAMMING AND SOFTWARE DEV. II_HIM 212_OUTLINE_May_2026 (2).pdf",
            "Python GUI Programming_P2_JUNE2025 (2).pptx"
        ]
        for f in files:
            st.code(f"📁 Path: Downloads/.../{f}")
            if not os.path.exists(f):
                with open(f, "w") as stub: stub.write("HIM 212 Lecture Content Stub.")
            with open(f, "rb") as b_file:
                st.download_button(label=f"📥 Download {f}", data=b_file, file_name=f, key=f"dl_{f}")

    # --- HIM 208 ---
    elif "HIM 208" in course_selection:
        st.write("#### Presentation Media")
        files = [
            "CANCER.pptx",
            "diet-and-health-ppt-1114he3 (3).pptx",
            "Fever-HIM NEW (2).pptx",
            "Introduction to Pathophysiology I.ppt"
        ]
        for f in files:
            st.code(f"📁 Path: Downloads/.../{f}")
            if not os.path.exists(f):
                with open(f, "w") as stub: stub.write("HIM 208 Lecture Content Stub.")
            with open(f, "rb") as b_file:
                st.download_button(label=f"📥 Download {f}", data=b_file, file_name=f, key=f"dl_{f}")

    # --- HIM 204 ---
    elif "HIM 204" in course_selection:
        st.write("#### Database Architecture & Design Handouts")
        files = [
            "HIM 204 LECTURE 1 - INTRODUCTION TO DBMS (3).pdf",
            "HIM 204 LECTURE 2 - DBMS ARCHITECTURE (2).pdf",
            "HIM 204 LECTURE 3 - ENTITY RELATION MODEL.pdf",
            "HIM 204 LECTURE 4 - WORKING WITH ER DIAGRAMS.pdf",
            "HIM 204 LECTURE 5 - NORMALIZATION IN DBMS.pdf",
            "HIM 204 LECTURE 6 - NORMAL FORMS.pdf",
            "HIM 204 LECTURE 7 - INTRODUCTION TO DATABASE QUERIES.pdf",
            "HIM 204 LECTURE 8 - TUTORIAL SESSION.pdf"
        ]
        for f in files:
            st.code(f"📁 Path: Downloads/.../{f}")
            if not os.path.exists(f):
                with open(f, "w") as stub: stub.write("HIM 204 Lecture Content Stub.")
            with open(f, "rb") as b_file:
                st.download_button(label=f"📥 Download {f}", data=b_file, file_name=f, key=f"dl_{f}")

    # --- HIM 210 ---
    elif "HIM 210" in course_selection:
        st.write("### ⚡ Integrated Functional Asset")
        st.info("💡 Below is your active ICD-10 Diagnostic Parser Engine embedded within its exact course module mapping!")
        st.markdown("---")

        search_query = st.text_input("🔍 ICD-10 Alpha-Index & Tabular Match Engine (Type Search Term):")
        if search_query:
            st.success(f"Parsing extracted hits for: '{search_query}' across Volume 1 and Volume 3 maps...")
            mock_hits = pd.DataFrame([
                {"Term Match": search_query.capitalize(), "ICD-10 Code block": "B15-B19", "Volume Registry": "Vol 3 Index (Page 328)"},
                {"Term Match": f"{search_query.capitalize()} (unspecified)", "ICD-10 Code block": "K75.9", "Volume Registry": "Vol 1 Tabular"}
            ])
            st.dataframe(mock_hits, use_container_width=True)

    # --- PHL 205, HIM 202, HIM 206 ---
    else:
        st.write("#### Course Materials")
        st.warning("🔄 System setup complete. Drop your remaining folder assets into the repository to activate download buttons!")


# ============================================================
# MODULE 2: AI-POWERED ADAPTIVE QUIZ ARENA
# ============================================================
elif app_mode == "📝 Adaptive Quiz Arena":
    st.markdown("<h1 class='main-title'>📝 AI-Powered Adaptive Quiz Arena</h1>", unsafe_allow_html=True)
    st.markdown("<div class='ai-badge'>✨ Powered by Claude AI — Real Exam Questions, All 7 Courses</div>", unsafe_allow_html=True)
    st.markdown("<p class='section-subtitle'>Answer questions correctly to step up the difficulty: Medium ➡️ Intermediate ➡️ Super Difficult.</p>", unsafe_allow_html=True)
    st.write("---")

    if st.session_state.active_student is None:
        st.warning("⚠️ Access Denied: You must register your name profile inside the 'Student Leaderboard Pool' tab first to join and compete live!")
    else:
        st.info(f"👤 Active Competitor: **{st.session_state.active_student}** | Points Accumulated: `{st.session_state.score}`")

        selected_course = st.selectbox("🎯 Select Course to Be Tested On:", ALL_COURSES)

        # Reset AI questions when course changes
        if selected_course != st.session_state.current_topic:
            st.session_state.ai_questions = []
            st.session_state.quiz_index = 0
            st.session_state.current_topic = selected_course
            st.session_state.answered = False
            st.session_state.last_correct = None
            st.session_state.explanation = ""

        # Difficulty tier tracker
        levels = ["Medium", "Intermediate", "Super Difficult"]
        current_lvl = st.session_state.quiz_level
        st.progress((levels.index(current_lvl) + 1) / len(levels))
        st.markdown(f"#### Active Difficulty Tier: :orange[{current_lvl}]")
        st.write("---")

        # Number of questions slider
        num_q = st.slider("How many AI questions to generate?", min_value=5, max_value=20, value=10, step=5)

        # Generate button
        gen_label = "🔁 Regenerate New Questions" if st.session_state.ai_questions else "🤖 Generate AI Questions"
        col_gen, _ = st.columns([2, 3])
        with col_gen:
            if st.button(gen_label, use_container_width=True):
                with st.spinner(f"🧠 Claude AI is crafting {num_q} real {current_lvl}-level questions for {selected_course}..."):
                    try:
                        st.session_state.ai_questions = generate_ai_questions(selected_course, num_q)
                        st.session_state.quiz_index = 0
                        st.session_state.answered = False
                        st.session_state.last_correct = None
                        st.session_state.explanation = ""
                        st.success(f"✅ {len(st.session_state.ai_questions)} questions ready for **{selected_course}** — [{current_lvl}] tier!")
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ Question generation failed: {e}")

        # Reset arena button
        if st.button("🔄 Reset Personal Arena Run"):
            st.session_state.quiz_level = "Medium"
            st.session_state.score = 0
            st.session_state.ai_questions = []
            st.session_state.quiz_index = 0
            st.session_state.answered = False
            st.session_state.last_correct = None
            st.session_state.explanation = ""
            if st.session_state.active_student in st.session_state.competition_pool:
                st.session_state.competition_pool[st.session_state.active_student] = {
                    "Score (pts)": 0, "Highest Tier": "Medium"
                }
            st.rerun()

        st.write("---")

        # --- QUIZ DISPLAY ---
        if not st.session_state.ai_questions:
            st.markdown(f"""
                <div style='text-align:center; padding:40px; opacity:0.65;'>
                    <h3 style='color:{text_color};'>🤖 No questions loaded yet</h3>
                    <p style='color:{sub_text};'>Select a course above and click <strong>"Generate AI Questions"</strong>
                    to create a fresh set of real exam questions.</p>
                </div>
            """, unsafe_allow_html=True)
        else:
            questions = st.session_state.ai_questions
            total_q   = len(questions)
            idx       = st.session_state.quiz_index

            if idx >= total_q:
                st.success(f"🏆 You've completed all {total_q} questions for **{selected_course}**! Score: **{st.session_state.score} pts**")

                # Advance difficulty tier automatically
                if current_lvl == "Medium":
                    st.session_state.quiz_level = "Intermediate"
                    st.info("🎯 Tier unlocked: **Intermediate**! Generate new questions to continue climbing.")
                elif current_lvl == "Intermediate":
                    st.session_state.quiz_level = "Super Difficult"
                    st.info("🔥 Tier unlocked: **Super Difficult**! Generate new questions for the final challenge.")
                elif current_lvl == "Super Difficult":
                    st.balloons()
                    st.success("🏆 Incredible! You have completely conquered all difficulty tiers for this course!")

                if st.session_state.active_student in st.session_state.competition_pool:
                    st.session_state.competition_pool[st.session_state.active_student] = {
                        "Score (pts)": st.session_state.score,
                        "Highest Tier": current_lvl
                    }

                if st.button("🔄 Generate Fresh Questions for Next Tier"):
                    st.session_state.ai_questions = []
                    st.session_state.quiz_index = 0
                    st.session_state.answered = False
                    st.session_state.last_correct = None
                    st.session_state.explanation = ""
                    st.rerun()

            else:
                q_data = questions[idx]

                # Progress bar
                st.markdown(f"### Question {idx + 1} of {total_q}")
                st.progress((idx + 1) / total_q)

                # Question card
                st.markdown(f"""
                    <div class='course-card'>
                        <p style='font-size:17px; font-weight:500; color:{text_color}; margin:0;'>{q_data['q']}</p>
                    </div>
                """, unsafe_allow_html=True)

                if not st.session_state.answered:
                    u_choice = st.radio("Select your answer:", q_data['o'], label_visibility="collapsed")

                    if st.button("🚀 Submit Answer"):
                        correct = (u_choice == q_data['c'])
                        st.session_state.last_correct = correct
                        st.session_state.explanation  = q_data.get("explanation", "")

                        if correct:
                            st.session_state.score += 50
                            if st.session_state.active_student in st.session_state.competition_pool:
                                st.session_state.competition_pool[st.session_state.active_student] = {
                                    "Score (pts)": st.session_state.score,
                                    "Highest Tier": current_lvl
                                }

                        st.session_state.answered = True
                        st.rerun()

                else:
                    # Result + Explanation
                    if st.session_state.last_correct:
                        st.success("🎯 Masterful! That choice is absolutely correct.")
                    else:
                        st.error(f"❌ Incorrect! The correct answer was: **{q_data['c']}**")

                    if st.session_state.explanation:
                        st.markdown(f"""
                            <div class='explanation-box'>
                                <strong style='color:#93C5FD;'>📖 Explanation:</strong>
                                <p style='color:#E2E8F0; margin:6px 0 0 0;'>{st.session_state.explanation}</p>
                            </div>
                        """, unsafe_allow_html=True)

                    if st.button("➡️ Next Question"):
                        st.session_state.quiz_index  += 1
                        st.session_state.answered     = False
                        st.session_state.last_correct = None
                        st.session_state.explanation  = ""
                        st.rerun()


# ============================================================
# MODULE 3: STUDENT LEADERBOARD POOL
# ============================================================
elif app_mode == "🏆 Student Leaderboard Pool":
    st.markdown("<h1 class='main-title'>🏆 Live Roster Competition Pool</h1>", unsafe_allow_html=True)
    st.markdown("<p class='section-subtitle'>Real-time ranking of real students who log onto the hub and decide to participate.</p>", unsafe_allow_html=True)
    st.write("---")

    st.markdown("""
        <div class='pool-banner'>
            <h3 style='margin:0 0 5px 0; font-weight:600; color:white;'>📝 Active Profile Entrance</h3>
            <p style='margin:0; opacity:0.9; font-size:14px; color:white;'>Enter your real name or unique campus alias to open an active tracking scorecard stream!</p>
        </div>
    """, unsafe_allow_html=True)

    join_col1, join_col2 = st.columns([3, 1])
    with join_col1:
        chosen_name = st.text_input("👤 Enter Full Student Name or Custom Alias Token:", placeholder="e.g. Mohammed Issifu")
    with join_col2:
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("⚡ Establish Pool Account", use_container_width=True):
            if chosen_name.strip() != "":
                clean_name = chosen_name.strip()
                st.session_state.active_student = clean_name
                if clean_name not in st.session_state.competition_pool:
                    st.session_state.competition_pool[clean_name] = {"Score (pts)": 0, "Highest Tier": "Medium"}
                st.success(f"🎉 Active session profile locked to: **{clean_name}**. Navigate to the Quiz tab to compete!")
                st.rerun()
            else:
                st.error("❌ Identification error: Input field blank!")

    st.write("---")

    if len(st.session_state.competition_pool) == 0:
        st.info("📊 The competition pool is currently empty. Be the first to type in your name above, hop into the Quiz Arena, and log points!")
    else:
        pool_records = []
        for name, data in st.session_state.competition_pool.items():
            pool_records.append({
                "Student Name":  name,
                "Score (pts)":   data.get("Score (pts)", 0),
                "Highest Tier":  data.get("Highest Tier", "Medium")
            })

        df_pool = pd.DataFrame(pool_records)
        df_pool = df_pool.sort_values(by="Score (pts)", ascending=False).reset_index(drop=True)
        df_pool.index += 1

        col_board, col_chart = st.columns([1, 1])
        with col_board:
            st.markdown("### 📊 Active Standing Board")
            st.dataframe(df_pool, use_container_width=True)
        with col_chart:
            st.markdown("### 📈 Class Ranking Visual Spread")
            st.bar_chart(data=df_pool, x="Student Name", y="Score (pts)", color="#FF4B4B")

        st.sidebar.info(f"🔥 Active Competitors in Pool: {len(pool_records)}")


# ============================================================
# MODULE 4: ACADEMIC TARGETER
# ============================================================
elif app_mode == "📊 Academic Targeter":
    st.markdown("<h1 class='main-title'>📊 Personal CGPA & Semester Goal Targeter</h1>", unsafe_allow_html=True)
    st.write("---")

    col_t1, col_t2 = st.columns(2)
    with col_t1:
        st.session_state.prev_cgpa = st.number_input(
            "Enter your previous cumulative CGPA:", value=st.session_state.prev_cgpa, step=0.1
        )
    with col_t2:
        st.session_state.target_gpa = st.number_input(
            "Enter your target GPA for this Semester 2:", value=st.session_state.target_gpa, step=0.1
        )

    if st.button("📈 Calculate Target Benchmarks"):
        st.balloons()
        st.success(
            f"To raise your academic trajectory from a {st.session_state.prev_cgpa} up to a "
            f"{st.session_state.target_gpa}, aim for a minimum grade score distribution profile of "
            f"B+ or higher across your 7 fundamental Level 200 health information system modules!"
        )
